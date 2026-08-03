"""PPTXファイルを解析し、各スライドの構成(タイトル/本文/ノート/画像有無など)を
構造化されたアウトライン(JSON + Markdown)として書き出す。

パターン2(既存PPTXから動画を作る)の最初のステップとして使う。
ここで出力したアウトラインを、チャットAI(docs/narration_prompt_template.md 参照)に
読ませて narration.md (ナレーション原稿) と index.html のスライド構成案を
生成してもらう、という想定。

実行例:
  python3 engine/pptx_outline.py projects/<project名>/source/foo.pptx
  -> projects/<project名>/source/foo.outline.json
     projects/<project名>/source/foo.outline.md
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def emu_to_px(emu, dpi=96):
    """EMU(English Metric Units)をおおよそのピクセル値に変換する。"""
    if emu is None:
        return None
    return round(Emu(emu).inches * dpi, 1)


def extract_text_frame(text_frame):
    """テキストフレームから段落ごとのテキストとインデントレベルを抽出する。"""
    paragraphs = []
    for para in text_frame.paragraphs:
        text = "".join(run.text for run in para.runs).strip()
        if not text:
            continue
        paragraphs.append({"level": para.level, "text": text})
    return paragraphs


def extract_shape(shape):
    info = {
        "shape_type": str(shape.shape_type) if shape.shape_type is not None else None,
        "name": shape.name,
        "position": {
            "left": emu_to_px(shape.left),
            "top": emu_to_px(shape.top),
            "width": emu_to_px(shape.width),
            "height": emu_to_px(shape.height),
        },
    }
    if shape.has_text_frame:
        info["text"] = extract_text_frame(shape.text_frame)
    if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
        info["is_picture"] = True
    if shape.has_table:
        table = shape.table
        info["table"] = [
            [cell.text.strip() for cell in row.cells] for row in table.rows
        ]
    return info


def extract_slide(slide, index):
    title = None
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        title = slide.shapes.title.text_frame.text.strip()

    body_paragraphs = []
    shapes_info = []
    has_picture = False
    for shape in slide.shapes:
        if shape == slide.shapes.title:
            continue
        shapes_info.append(extract_shape(shape))
        if shape.has_text_frame:
            body_paragraphs.extend(extract_text_frame(shape.text_frame))
        if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
            has_picture = True

    notes = ""
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()

    layout_name = slide.slide_layout.name if slide.slide_layout else None

    return {
        "index": index,
        "layout": layout_name,
        "title": title,
        "body": body_paragraphs,
        "shapes": shapes_info,
        "has_picture": has_picture,
        "notes": notes,
    }


def build_outline(pptx_path):
    prs = Presentation(str(pptx_path))
    slides = [extract_slide(slide, i + 1) for i, slide in enumerate(prs.slides)]
    return {
        "source": str(pptx_path),
        "slide_size": {
            "width": emu_to_px(prs.slide_width),
            "height": emu_to_px(prs.slide_height),
        },
        "slide_count": len(slides),
        "slides": slides,
    }


def render_markdown(outline):
    lines = [f"# PPTXアウトライン: {Path(outline['source']).name}", ""]
    lines.append(f"- スライド数: {outline['slide_count']}")
    lines.append("")
    for slide in outline["slides"]:
        lines.append(f"## Slide {slide['index']}" + (f": {slide['title']}" if slide["title"] else ""))
        if slide["layout"]:
            lines.append(f"- レイアウト: {slide['layout']}")
        if slide["has_picture"]:
            lines.append("- 画像あり")
        if slide["body"]:
            lines.append("- 本文:")
            for p in slide["body"]:
                indent = "  " * (p["level"] + 1)
                lines.append(f"{indent}- {p['text']}")
        if slide["notes"]:
            lines.append(f"- ノート: {slide['notes']}")
        lines.append("")
    return "\n".join(lines)


def main():
    if len(sys.argv) < 2:
        print("使い方: pptx_outline.py <pptxファイルのパス>", file=sys.stderr)
        raise SystemExit(1)

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        raise SystemExit(f"ファイルが見つかりません: {pptx_path}")

    outline = build_outline(pptx_path)

    json_path = pptx_path.with_suffix(".outline.json")
    md_path = pptx_path.with_suffix(".outline.md")

    json_path.write_text(
        json.dumps(outline, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    md_path.write_text(render_markdown(outline), encoding="utf-8")

    print(f"スライド数: {outline['slide_count']}")
    print(f"-> {json_path}")
    print(f"-> {md_path}")


if __name__ == "__main__":
    main()
