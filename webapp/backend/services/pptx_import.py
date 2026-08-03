"""PPTX を決定論的に解析し、1 スライド = 1 シーンとして構造化する。

LLM には一切依存しない。シーン分割・レイアウト決定はすべてルールベースで行い、
LLM の役割はナレーション文の生成（services/llm_service.py 側）だけに限定する。

FIX-15 (docs/antigravity_fix15_pptx_per_slide_import.md) の解析部。
"""
from __future__ import annotations

import hashlib
import io
import sys
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from services.pptx_vector import (
    is_available as vector_is_available,
    pptx_to_page_svgs,
    crop_svg,
)

MIN_AREA_RATIO = 0.015          # スライド面積の1.5%未満のビジュアルは無視（装飾・アイコン）
MAX_VISUALS_PER_SLIDE = 3       # SceneAsset の slot は 1〜3
# 取り込み画像の長辺の上限。動画は 1920x1080 なのでこれを超える解像度は無駄になる。
# PowerPoint には原寸のまま数MBの画像が貼られていることが多く、全シーン分が
# 1ページに載るレンダリング時に Chrome のメモリを大きく圧迫するため縮小する。
MAX_IMAGE_LONG_SIDE = 1920
JPEG_QUALITY = 85
LOGO_MIN_COUNT = 3               # 同一バイナリがこの枚数以上、かつ
LOGO_MIN_RATIO = 0.6             # 全スライドの60%以上に出現するものはロゴ扱いで除外
MAX_SVG_BYTES = 2 * 1024 * 1024  # SVG が2MBを超える場合は採用しない


@dataclass
class SlideVisual:
    kind: str                       # "raster" | "vector"
    ext: str                        # ".png" / ".jpg" / ".svg" など
    data: bytes
    bbox_emu: tuple[int, int, int, int]   # left, top, width, height
    area_ratio: float
    digest: str
    chart_data: dict | None = None


@dataclass
class SlideInfo:
    index: int                      # 1始まり。表示スライドのみの連番
    title: str = ""
    bullets: list[str] = field(default_factory=list)
    body_text: str = ""
    notes: str = ""
    table: dict | None = None       # {"headers": [...], "rows": [[...]]}
    visuals: list[SlideVisual] = field(default_factory=list)
    layout_type: str = "text_only"
    image_position: str = "right"   # "left" | "right"
    warnings: list[str] = field(default_factory=list)


def _emu_area(w, h) -> int:
    return int(w or 0) * int(h or 0)


def _is_title_shape(slide, shape) -> bool:
    title_shape = slide.shapes.title
    if title_shape is not None and shape.shape_id == title_shape.shape_id:
        return True
    if getattr(shape, "is_placeholder", False):
        ph = shape.placeholder_format
        if ph is not None and ph.type is not None and str(ph.type) in ("TITLE (13)", "CENTER_TITLE (0)"):
            return True
    return False


def _is_smartart(shape) -> bool:
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return False
        el = shape._element
        graphic_data = el.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}graphicData")
        if graphic_data is not None:
            uri = graphic_data.get("uri", "")
            return "diagram" in uri
    except Exception:
        pass
    return False


def optimize_raster(data: bytes, ext: str) -> tuple[bytes, str]:
    """取り込んだ画像を動画用に最適化する（長辺 1920px 上限・不要なアルファは JPEG 化）。

    戻り値: (最適化後のバイト列, 拡張子)。失敗した場合は元データをそのまま返す。
    """
    try:
        from PIL import Image
    except ImportError:
        return data, ext

    try:
        with Image.open(io.BytesIO(data)) as im:
            im.load()
            # アニメーション GIF は加工すると壊れるのでそのまま使う
            if getattr(im, "is_animated", False):
                return data, ext

            has_alpha = im.mode in ("RGBA", "LA", "PA") or "transparency" in im.info
            long_side = max(im.width, im.height)
            needs_resize = long_side > MAX_IMAGE_LONG_SIDE

            if not needs_resize and (has_alpha or ext in (".jpg", ".jpeg")):
                # 縮小不要で、かつ再エンコードしても得がない場合は元のまま
                return data, ext

            if needs_resize:
                scale = MAX_IMAGE_LONG_SIDE / long_side
                im = im.resize(
                    (max(1, round(im.width * scale)), max(1, round(im.height * scale))),
                    Image.LANCZOS,
                )

            buf = io.BytesIO()
            if has_alpha:
                # 透過を保つ必要があるので PNG のまま
                im.save(buf, format="PNG", optimize=True)
                out_ext = ".png"
            else:
                im.convert("RGB").save(buf, format="JPEG", quality=JPEG_QUALITY, optimize=True)
                out_ext = ".jpg"

            optimized = buf.getvalue()
            # 最適化して大きくなるなら元を採用する
            if len(optimized) >= len(data) and not needs_resize:
                return data, ext
            return optimized, out_ext
    except Exception:
        return data, ext


def _extract_visual_from_picture(shape, slide_area: int) -> SlideVisual | None:
    try:
        image = shape.image
    except Exception:
        return None
    data = image.blob
    ext = "." + (image.ext or "png").lower()
    if ext not in (".png", ".jpg", ".jpeg", ".gif", ".webp"):
        # EMF/WMF 等のベクタ形式はラスタとしては扱わない（呼び出し元で図解クラスタに回す）
        return None
    # レンダリング時のメモリ削減のため、ここで縮小・再エンコードしておく
    data, ext = optimize_raster(data, ext)
    area = _emu_area(shape.width, shape.height)
    ratio = area / slide_area if slide_area else 0
    digest = hashlib.sha1(data).hexdigest()
    return SlideVisual(
        kind="raster",
        ext=ext if ext != ".jpeg" else ".jpg",
        data=data,
        bbox_emu=(int(shape.left or 0), int(shape.top or 0), int(shape.width or 0), int(shape.height or 0)),
        area_ratio=ratio,
        digest=digest,
    )


def _extract_chart_data(shape) -> dict | None:
    try:
        chart = shape.chart
        chart_type = "bar"
        type_name = str(chart.chart_type)
        if "LINE" in type_name:
            chart_type = "line"
        elif "PIE" in type_name:
            chart_type = "pie"
        plot = chart.plots[0]
        labels = [str(c) for c in plot.categories]
        series = list(chart.series)
        values = [float(v) if v is not None else 0.0 for v in series[0].values] if series else []
        return {"type": chart_type, "labels": labels, "values": values, "unit": ""}
    except Exception:
        return None


def _bbox_union(bboxes: list[tuple[int, int, int, int]]) -> tuple[int, int, int, int]:
    lefts = [b[0] for b in bboxes]
    tops = [b[1] for b in bboxes]
    rights = [b[0] + b[2] for b in bboxes]
    bottoms = [b[1] + b[3] for b in bboxes]
    l, t = min(lefts), min(tops)
    return (l, t, max(rights) - l, max(bottoms) - t)


def _collect_diagram_shapes(slide, slide_area: int):
    """ラスタ画像でもテキストだけのシェイプでもない、図解候補シェイプを集める。

    戻り値: [(kind, shape_or_group_bboxes)] のリスト。
      - "chart" / "smartart" / "group" は単独で1クラスタ
      - それ以外(オートシェイプ/コネクタ/フリーフォーム/EMF画像)はスライドで1クラスタにまとめる
    """
    chart_clusters = []
    smartart_clusters = []
    group_clusters = []
    misc_bboxes = []
    misc_has_only_text = True

    for shape in slide.shapes:
        st = shape.shape_type
        try:
            if st == MSO_SHAPE_TYPE.CHART or getattr(shape, "has_chart", False):
                chart_clusters.append(shape)
                continue
        except Exception:
            pass
        if _is_smartart(shape):
            smartart_clusters.append(shape)
            continue
        if st == MSO_SHAPE_TYPE.GROUP:
            group_clusters.append(shape)
            continue
        if st == MSO_SHAPE_TYPE.PICTURE:
            try:
                ext = ("." + (shape.image.ext or "")).lower()
            except Exception:
                ext = ""
            if ext not in (".png", ".jpg", ".jpeg", ".gif", ".webp"):
                # EMF/WMF など: 図解候補として集める
                misc_bboxes.append((int(shape.left or 0), int(shape.top or 0), int(shape.width or 0), int(shape.height or 0)))
                misc_has_only_text = False
            continue
        if st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.LINE) or shape.shape_type is None:
            if shape.left is None or shape.width is None:
                continue
            misc_bboxes.append((int(shape.left or 0), int(shape.top or 0), int(shape.width or 0), int(shape.height or 0)))
            if not shape.has_text_frame or not shape.text_frame.text.strip():
                misc_has_only_text = False

    clusters = []
    for s in chart_clusters:
        clusters.append(("chart", s, (int(s.left or 0), int(s.top or 0), int(s.width or 0), int(s.height or 0))))
    for s in smartart_clusters:
        clusters.append(("smartart", s, (int(s.left or 0), int(s.top or 0), int(s.width or 0), int(s.height or 0))))
    for s in group_clusters:
        clusters.append(("group", s, (int(s.left or 0), int(s.top or 0), int(s.width or 0), int(s.height or 0))))

    # misc は単独のテキストシェイプだけなら図解とみなさない（2個以上、または非テキスト要素が混じる場合のみ採用）
    if misc_bboxes and (len(misc_bboxes) >= 2 or not misc_has_only_text):
        bbox = _bbox_union(misc_bboxes)
        clusters.append(("misc", None, bbox))

    result = []
    for kind, shape, bbox in clusters:
        area = bbox[2] * bbox[3]
        ratio = area / slide_area if slide_area else 0
        if ratio < MIN_AREA_RATIO:
            continue
        result.append((kind, shape, bbox))
    return result


def _decide_layout(visual_count: int, has_table: bool, bullets: list[str], body_text: str, title: str, layout_name: str) -> str:
    if visual_count >= 2:
        return "image_gallery"
    if visual_count == 1 and len(body_text.strip()) < 20 and not bullets:
        return "full_image"
    if visual_count == 1:
        return "text_left_image_right"
    if has_table:
        return "table"
    if len(bullets) >= 2:
        return "bullet_list"
    has_body = bool(body_text.strip()) or len(bullets) > 0
    if not has_body:
        name = (layout_name or "").lower()
        if "title" in name or "タイトル" in (layout_name or ""):
            return "section_header"
        if title and not body_text and not bullets:
            return "section_header"
    return "text_only"


def parse_pptx(pptx_path: Path, workdir: Path) -> tuple[list[SlideInfo], list[str]]:
    """PPTX を解析し、(スライド情報のリスト, 全体警告リスト) を返す。

    スライド情報のリストの長さは、必ず「表示スライド数」と一致する。
    """
    warnings: list[str] = []
    prs = Presentation(str(pptx_path))
    slide_w, slide_h = int(prs.slide_width), int(prs.slide_height)
    slide_area = slide_w * slide_h

    visible_slides = []
    for slide in prs.slides:
        show = slide.element.get("show")
        if show == "0":
            continue
        visible_slides.append(slide)

    if not visible_slides:
        return [], ["有効なスライドが見つかりませんでした"]

    use_vector = vector_is_available()
    page_svgs: dict[int, str] = {}
    if use_vector:
        try:
            page_svgs = pptx_to_page_svgs(pptx_path, workdir, len(visible_slides))
            if len(page_svgs) != len(visible_slides):
                warnings.append("スライド枚数とSVGページ数が一致しないため、図解のSVG化を一部スキップしました")
        except Exception as e:
            warnings.append(f"LibreOffice によるSVG変換に失敗しました: {e}")
            page_svgs = {}
    else:
        warnings.append("図解をSVG化できませんでした（LibreOffice未導入）。図解はテキスト箇条書きとして取り込みます")

    # ロゴ判定用: 全スライドを通じたラスタ画像の出現回数
    digest_counts: dict[str, int] = {}

    slides_info: list[SlideInfo] = []
    for i, slide in enumerate(visible_slides, start=1):
        info = SlideInfo(index=i)

        title_shape = slide.shapes.title
        if title_shape is not None and title_shape.has_text_frame:
            info.title = title_shape.text_frame.text.strip()

        # テーブルは別途抽出。段落が複数あるテキストシェイプは箇条書き、1段落だけなら本文として扱う
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
                if rows:
                    info.table = {"headers": rows[0], "rows": rows[1:]}

        bullets: list[str] = []
        body_lines: list[str] = []
        for shape in slide.shapes:
            if _is_title_shape(slide, shape) or not shape.has_text_frame:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue
            paras = [p for p in shape.text_frame.paragraphs]
            texts = []
            for p in paras:
                t = "".join(r.text for r in p.runs).strip()
                if t:
                    texts.append(t)
            if not texts:
                continue
            if len(texts) > 1:
                bullets.extend(texts)
            else:
                body_lines.append(texts[0])

        info.bullets = bullets
        info.body_text = "\n".join(body_lines)

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            info.notes = slide.notes_slide.notes_text_frame.text.strip()

        # ── ビジュアル抽出 ──
        raw_visuals: list[SlideVisual] = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                v = _extract_visual_from_picture(shape, slide_area)
                if v:
                    raw_visuals.append(v)
                    digest_counts[v.digest] = digest_counts.get(v.digest, 0) + 1

        diagram_clusters = _collect_diagram_shapes(slide, slide_area)
        page_svg = page_svgs.get(i)
        for kind, shape, bbox in diagram_clusters:
            chart_data = None
            if kind == "chart" and shape is not None:
                chart_data = _extract_chart_data(shape)
            if page_svg:
                try:
                    svg_str = crop_svg(page_svg, bbox, (slide_w, slide_h))
                    svg_bytes = svg_str.encode("utf-8")
                    if len(svg_bytes) <= MAX_SVG_BYTES:
                        area = bbox[2] * bbox[3]
                        ratio = area / slide_area if slide_area else 0
                        raw_visuals.append(SlideVisual(
                            kind="vector", ext=".svg", data=svg_bytes,
                            bbox_emu=bbox, area_ratio=ratio,
                            digest=hashlib.sha1(svg_bytes).hexdigest(),
                            chart_data=chart_data,
                        ))
                        continue
                    else:
                        info.warnings.append("図解のSVGサイズが大きすぎるため取り込みをスキップしました")
                except Exception as e:
                    info.warnings.append(f"図解のSVG切り出しに失敗しました: {e}")
            # SVG化できない場合のフォールバック
            if kind == "chart" and chart_data and chart_data.get("labels"):
                labels = chart_data["labels"]
                values = chart_data["values"]
                bullets.extend([f"{l}: {v}" for l, v in zip(labels, values)])
            elif shape is not None and shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    bullets.append(t)

        raw_visuals.sort(key=lambda v: v.area_ratio, reverse=True)
        info.visuals = raw_visuals[:MAX_VISUALS_PER_SLIDE]

        if info.visuals:
            first = info.visuals[0]
            center_x = first.bbox_emu[0] + first.bbox_emu[2] / 2
            info.image_position = "left" if center_x < slide_w / 2 else "right"

        info.layout_type = _decide_layout(
            visual_count=len(info.visuals),
            has_table=info.table is not None,
            bullets=info.bullets,
            body_text=info.body_text,
            title=info.title,
            layout_name=slide.slide_layout.name if slide.slide_layout else "",
        )

        slides_info.append(info)

    # ロゴ判定: 出現回数が閾値以上のダイジェストを持つビジュアルを除外
    total_slides = len(slides_info)
    logo_digests = {
        d for d, c in digest_counts.items()
        if c >= LOGO_MIN_COUNT and (c / total_slides) >= LOGO_MIN_RATIO
    }
    if logo_digests:
        for info in slides_info:
            before = len(info.visuals)
            info.visuals = [v for v in info.visuals if v.digest not in logo_digests]
            if len(info.visuals) != before:
                # レイアウトを再決定（ビジュアルが減った可能性があるため）
                info.layout_type = _decide_layout(
                    visual_count=len(info.visuals),
                    has_table=info.table is not None,
                    bullets=info.bullets,
                    body_text=info.body_text,
                    title=info.title,
                    layout_name="",
                ) if info.layout_type in ("image_gallery", "full_image", "text_left_image_right") else info.layout_type

    return slides_info, warnings


def _cli():
    if len(sys.argv) < 2:
        print("使い方: python3 -m services.pptx_import <pptxファイル> [--out DIR]", file=sys.stderr)
        raise SystemExit(1)
    pptx_path = Path(sys.argv[1])
    out_dir = Path("/tmp/pptx_check")
    if "--out" in sys.argv:
        out_dir = Path(sys.argv[sys.argv.index("--out") + 1])
    out_dir.mkdir(parents=True, exist_ok=True)

    slides, warnings = parse_pptx(pptx_path, out_dir)
    for w in warnings:
        print(f"[WARN] {w}", file=sys.stderr)
    for s in slides:
        kinds = ",".join(v.kind for v in s.visuals) or "-"
        print(f"slide {s.index:>2}: layout={s.layout_type:<22} title={s.title[:20]!r:<24} "
              f"visuals={len(s.visuals)}({kinds}) bullets={len(s.bullets)} table={'Y' if s.table else 'N'}")
        for j, v in enumerate(s.visuals, start=1):
            fp = out_dir / f"slide{s.index}_slot{j}{v.ext}"
            fp.write_bytes(v.data)
            print(f"    -> {fp}")


if __name__ == "__main__":
    _cli()
